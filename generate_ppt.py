import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
import re

def create_ppt():
    with open('presentation.md', 'r', encoding='utf-8') as f:
        content = f.read()

    slides_data = content.split('## ')
    prs = Presentation()

    for slide_raw in slides_data[1:]:
        lines = slide_raw.strip().split('\n')
        title_line = lines[0].strip()
        # Clean title by removing slide prefix and emojis
        title = re.sub(r'Slide \d+: ', '', title_line)
        title = re.sub(r'[^\w\s\-\:\.\,]', '', title).strip() # basic emoji removal
        
        body_lines = []
        for line in lines[1:]:
            line = line.strip()
            if line and line != '---':
                body_lines.append(line)
                
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        first = True
        for line in body_lines:
            clean_line = line.replace('**', '').replace('`', '')
            if clean_line.startswith('- '):
                clean_line = clean_line[2:]
            
            if first:
                tf.text = clean_line
                first = False
            else:
                p = tf.add_paragraph()
                p.text = clean_line

    prs.save('Phishield_AI_Presentation.pptx')
    print("Presentation saved as Phishield_AI_Presentation.pptx")

if __name__ == '__main__':
    create_ppt()
