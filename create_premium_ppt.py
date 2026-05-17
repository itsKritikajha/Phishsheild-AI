import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def set_dark_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 15, 25)  # Dark cyber blue/black

def add_title(slide, text, color=RGBColor(0, 240, 255)):
    title_shape = slide.shapes.title
    title_shape.text = text
    p = title_shape.text_frame.paragraphs[0]
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = 'Arial'

def add_bullets(slide, items):
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for item in items:
        p = tf.add_paragraph()
        if item.startswith("  - "):
            p.text = item[4:]
            p.level = 1
            p.font.size = Pt(18)
        elif item.startswith("- "):
            p.text = item[2:]
            p.level = 0
            p.font.size = Pt(22)
        else:
            p.text = item
            p.level = 0
            p.font.size = Pt(22)
            p.font.bold = True
        
        p.font.color.rgb = RGBColor(220, 230, 255)
        p.font.name = 'Arial'

def create_diagram(slide, steps, start_top=2.0):
    # Create a vertical flow diagram
    left = Inches(1.5)
    width = Inches(7)
    height = Inches(0.6)
    
    for i, step in enumerate(steps):
        top = Inches(start_top + (i * 0.8))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(20, 30, 50)
        shape.line.color.rgb = RGBColor(0, 240, 255)
        shape.line.width = Pt(1.5)
        
        tf = shape.text_frame
        tf.text = step
        p = tf.paragraphs[0]
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Add arrow pointing down (except last)
        if i < len(steps) - 1:
            arrow_top = top + height
            arrow_height = Inches(0.2)
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.8), arrow_top, Inches(0.4), arrow_height)
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(180, 50, 255) # Neon purple
            arrow.line.color.rgb = RGBColor(180, 50, 255)

def create_premium_ppt():
    prs = Presentation()
    
    # 1. Cover Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_dark_bg(slide)
    
    title = slide.shapes.title
    title.text = "Phishield AI"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 240, 255)
    title.text_frame.paragraphs[0].font.size = Pt(64)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle = slide.placeholders[1]
    tf = subtitle.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "AI Powered Phishing Detection & Cyber Threat Protection\n"
    p.font.color.rgb = RGBColor(180, 50, 255) # Cyber purple
    p.font.size = Pt(24)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "\nTeam:\nKritika Jha\nTarun Kushwaha\n\nDepartment: B.Tech CSE AIML"
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.font.size = Pt(20)

    # 2. Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)
    add_title(slide, "Problem Statement", RGBColor(255, 50, 100))
    add_bullets(slide, [
        "- Rise of phishing attacks worldwide",
        "- Fake websites and malicious links deceive users",
        "- Escalation in cyber fraud and identity theft",
        "- Severe lack of awareness among general users",
        "Modern Phishing Statistics:",
        "  - Over 3.4 billion malicious emails are sent daily",
        "  - Typosquatting URLs bypass traditional filters",
        "  - 90% of breaches involve human manipulation"
    ])

    # 3. Our Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)
    add_title(slide, "Our Solution: Phishield AI")
    add_bullets(slide, [
        "Phishield AI solves the problem using advanced technology:",
        "- Artificial Intelligence & Machine Learning",
        "- Real-time URL and semantic analysis",
        "- Proactive threat detection system",
        "- Instant alerts before users are compromised",
        "Key Differentiators:",
        "  - Analyzes behavioral patterns instead of just blacklists",
        "  - Calculates structural and lexical similarities",
        "  - Integrates directly into browser sessions"
    ])

    # 4. System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
    set_dark_bg(slide)
    add_title(slide, "System Architecture")
    create_diagram(slide, [
        "User Input URL",
        "Frontend Interface",
        "Feature Extraction Engine",
        "Machine Learning Model",
        "Threat Analysis Engine",
        "Prediction Output",
        "Dashboard & Alerts"
    ], start_top=1.5)

    # 5. Working Flow Diagram
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title only
    set_dark_bg(slide)
    add_title(slide, "Working Flow Diagram")
    create_diagram(slide, [
        "1. User enters suspicious URL",
        "2. System extracts URL features (Length, Special Chars, NLP)",
        "3. AI model analyzes patterns",
        "4. Threat score generated (0-100 Confidence)",
        "5. URL classified as: Safe / Suspicious / Dangerous",
        "6. Result visually shown on Dashboard",
        "7. Browser extension blocks/alerts user instantly"
    ], start_top=1.5)

    # 6. Tech Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)
    add_title(slide, "Technology Stack")
    add_bullets(slide, [
        "Frontend:",
        "  - HTML5, CSS3 (Glassmorphism & Neon UI), JavaScript",
        "Backend:",
        "  - Python (Flask Framework)",
        "AI / Machine Learning:",
        "  - Scikit-Learn, Pandas, Joblib",
        "  - Random Forest Classifier",
        "Browser Extension:",
        "  - Chrome Manifest V3 API (JavaScript)",
        "Database / Storage:",
        "  - LocalStorage & File System Logging"
    ])

    # 7. AI & Machine Learning
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)
    add_title(slide, "AI & Machine Learning")
    add_bullets(slide, [
        "Feature Extraction:",
        "  - URL Length, Subdomains, Hyphens, HTTP/HTTPS usage",
        "Algorithms Deployed:",
        "  - Random Forest Model (High Precision, Ensemble Learning)",
        "  - Logistic Regression (Fallback Analysis)",
        "Capabilities:",
        "  - Evaluates lexical and structural similarities",
        "  - Generates an AI Confidence Score",
        "  - Analyzes email semantics via NLP techniques"
    ])

    # 8. Browser Extension
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)
    add_title(slide, "Browser Extension Integration")
    add_bullets(slide, [
        "- Live phishing detection on active tabs",
        "- Dynamic Website Trust Score",
        "- Real-time popup warnings",
        "- Seamless Browser protection system",
        "  - Prevents accidental clicks on malicious payloads",
        "  - Communicates directly with the Python API Backend",
        "  - Lightweight, fast, and un-intrusive UI"
    ])

    # 9. Dashboard & Analytics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)
    add_title(slide, "Dashboard & Analytics")
    add_bullets(slide, [
        "- Global Threat Analytics Dashboard",
        "- Live AI detection stats & Threat Maps",
        "- Comprehensive Scan History",
        "- Interactive Risk Meter & AI Explainability Graphs",
        "- Enterprise Reporting:",
        "  - Automated PDF and CSV exports",
        "  - Real-time global threat feed simulation"
    ])

    # 10. Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)
    add_title(slide, "Key Features")
    add_bullets(slide, [
        "- Real-time URL Scanning",
        "- AI Phishing & Clone Site Detection",
        "- Multi-Level Threat Classification",
        "- Live Browser Extension",
        "- Advanced Dashboard Analytics",
        "- Proactive Cybersecurity Awareness via AI Assistant",
        "- Instantaneous Threat Alerts & Visual Gauges"
    ])

    # 11. Real World Impact
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)
    add_title(slide, "Real World Impact")
    add_bullets(slide, [
        "- User Protection: Shields individuals from zero-day phishing",
        "- Banking Safety: Prevents financial credentials from theft",
        "- Scam Prevention: Reduces success rates of social engineering",
        "- Cyber Awareness: Educates users with AI explanations",
        "- Enterprise Defense: Blocks attacks before they breach the network"
    ])

    # 12. Future Scope
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_dark_bg(slide)
    add_title(slide, "Future Scope")
    add_bullets(slide, [
        "- QR Phishing Detection (Quishing)",
        "- Full IMAP Email Phishing Scanner Integration",
        "- Autonomous Multi-Agent AI Chatbot",
        "- Dedicated Mobile App for link verification",
        "- Live Global Cyber Threat Map Integrations",
        "- Voice Assistant alerts for visually impaired users"
    ])

    # 13. Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_dark_bg(slide)
    title = slide.shapes.title
    title.text = "Conclusion"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 240, 255)
    
    subtitle = slide.placeholders[1]
    tf = subtitle.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "“Phishield AI aims to create a safer digital ecosystem using Artificial Intelligence and modern cybersecurity solutions.”"
    p.font.color.rgb = RGBColor(220, 230, 255)
    p.font.size = Pt(28)
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "\nThank You!\nQuestions?"
    p.font.color.rgb = RGBColor(180, 50, 255)
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    prs.save('Phishield_AI_Premium_Presentation.pptx')
    print("Premium presentation saved as Phishield_AI_Premium_Presentation.pptx")

if __name__ == '__main__':
    create_premium_ppt()
